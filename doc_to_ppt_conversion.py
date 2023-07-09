from pptx import Presentation
from pptx.util import Inches
from docx import Document
from pylatexenc.latex2text import LatexNodes2Text

def convert_docx_to_ppt(docx_file_path, ppt_file_path):
    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Load the DOCX file as a document object
    doc = Document(docx_file_path)

    # Iterate over the paragraphs in the document
    for paragraph in doc.paragraphs:
        # Add a new slide with a title and content layout
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the text content of the slide
        title = slide.shapes.title
        title.text = "Formula Slide"

        # Create a content placeholder on the slide
        content_placeholder = slide.placeholders[1]

        # Get the LaTeX formula from the paragraph
        latex_formula = paragraph.text

        # Convert the LaTeX formula to plain text
        formula_text = LatexNodes2Text().latex_to_text(latex_formula)

        # Add the formula text to the content placeholder
        content_placeholder.text = formula_text

    # Save the PowerPoint presentation
    presentation.save(ppt_file_path)

# Specify the paths for the input DOCX file and the output PPT file
docx_file_path = r"C:\Users\bhuth\Downloads\Sample_DOCX (1).docx"
ppt_file_path = r"C:\Users\bhuth\OneDrive\Desktop\vsb.pptx"

# Convert the DOCX file to a PowerPoint presentation
convert_docx_to_ppt(docx_file_path, ppt_file_path)
