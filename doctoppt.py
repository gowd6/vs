from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Pt

def convert_docx_to_ppt(docx_file_path, ppt_file_path):
    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Load the DOCX file as a document object
    doc = Document(docx_file_path)

    # Iterate over the paragraphs in the document
    for paragraph in doc.paragraphs:
        # Add a new slide with a title and content layout
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the text content of the slide
        title = slide.shapes.title
        title.text = "Formula Slide"

        # Create a content placeholder on the slide
        content_placeholder = slide.placeholders[1]

        # Iterate over the runs in the paragraph
        for run in paragraph.runs:
            # Check if the run has an inline picture
            if run.inline_shapes:
                # Add the picture to the content placeholder
                inline_shape = run.inline_shapes[0]
                content_placeholder.shapes.add_picture(
                    inline_shape.image.filename, Inches(2), Inches(2), width=Inches(4), height=Inches(2)
                )
            else:
                # Add the run text as text content to the content placeholder
                p = content_placeholder.text_frame.add_paragraph()
                p.text = run.text

                # Set the font size of the text content
                p.font.size = Pt(14)

    # Save the PowerPoint presentation
    presentation.save(ppt_file_path)

# Specify the paths for the input DOCX file and the output PPT file


# Specify the paths for the input DOCX file and the output PPT file
docx_file_path = r"C:\Users\bhuth\Downloads\Sample_DOCX (1).docx"
ppt_file_path = r"C:\Users\bhuth\OneDrive\Desktop\vsb.pptx"

# Convert the DOCX file to a PowerPoint presentation
convert_docx_to_ppt(docx_file_path, ppt_file_path)
